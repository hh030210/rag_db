import openpyxl
import pandas as pd
from pptx import Presentation
import pdfplumber
import docx
def parse_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"[错误] 解析Word文件失败: {e}")
        return f"解析Word文件失败: {e}"

def parse_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"[错误] 解析PowerPoint文件失败: {e}")
        return f"解析PowerPoint文件失败: {e}"

def parse_xlsx(file_path):
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        output = []
        for sheet_name, df in df_dict.items():
            output.append(f"Sheet: {sheet_name}")
            output.append(df.to_string())
        return "\n".join(output)
    except Exception as e:
        print(f"[错误] 解析Excel文件失败: {e}")
        return f"解析Excel文件失败: {e}"

def parse_editable_pdf(file_path):
    try:
        with pdfplumber.open(file_path) as pdf:
            content = []
            for page in pdf.pages:
                content.append(page.extract_text() or "")
        return "\n".join(content)
    except Exception as e:
        print(f"[错误] 解析PDF文件失败: {e}")
        return f"解析PDF文件失败: {e}"

def parse_editable(file_path, file_type):
    if file_type == 'docx':
        return parse_docx(file_path)
    elif file_type == 'pptx':
        return parse_pptx(file_path)
    elif file_type == 'xlsx':
        return parse_xlsx(file_path)
    elif file_type == 'pdf':
        return parse_editable_pdf(file_path)
    return ""
